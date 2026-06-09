"""Take screenshot of org chart HTML with white background, embed in PPTX."""
import os
import sys
import time
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image


def main(html_path, pptx_path):
    screenshot_path = os.path.join(os.path.dirname(pptx_path), "_screenshot.png")

    with sync_playwright() as p:
        try:
            browser = p.chromium.launch(channel="msedge")
        except Exception as e:
            print(f"ERROR: Cannot launch Edge. Details: {e}")
            sys.exit(1)

        page = browser.new_page(viewport={"width": 1400, "height": 4000})
        try:
            # Convert path to proper file URI for cross-platform compatibility
            file_uri = Path(html_path).absolute().as_uri()
            page.goto(file_uri, wait_until="networkidle", timeout=15000)
        except Exception as e:
            print(f"ERROR: Failed to load HTML: {e}")
            browser.close()
            sys.exit(1)

        try:
            # Wait for SVG lines to render
            page.wait_for_function(
                "() => { var svg = document.getElementById('lines'); return svg && svg.innerHTML.trim() !== ''; }",
                timeout=5000
            )
        except Exception:
            pass

        # Wait for fonts to load (critical for Chinese characters)
        page.wait_for_function("() => document.fonts.ready", timeout=3000)

        # Override body/html background to white
        page.evaluate("""() => {
            var sheet = document.createElement('style');
            sheet.textContent = 'html, body { background: #FFFFFF !important; }';
            document.head.appendChild(sheet);
        }""")

        # Force reflow to ensure fonts are applied
        page.evaluate("() => document.body.offsetHeight")

        # Additional wait for rendering
        page.wait_for_timeout(300)

        # Measure content dimensions dynamically
        dims = page.evaluate("""() => {
            var mc = document.getElementById('mainContainer');
            var cards = document.querySelectorAll('.nc');
            var maxRight = 0, maxBottom = 0;
            cards.forEach(function(c) {
                var r = c.getBoundingClientRect();
                if (r.right > maxRight) maxRight = r.right;
                if (r.bottom > maxBottom) maxBottom = r.bottom;
            });
            return {
                width: Math.max(maxRight + 40, 800),
                height: Math.max(maxBottom + 20, 400),
            };
        }""")

        clip_w = min(dims["width"], 4000)
        clip_h = min(dims["height"], 4000)
        page.screenshot(path=screenshot_path, clip={"x": 0, "y": 0, "width": clip_w, "height": clip_h})
        browser.close()

    print(f"Screenshot: {screenshot_path} ({os.path.getsize(screenshot_path)} bytes)")

    # Embed in PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    img = Image.open(screenshot_path)
    img_w_in = img.width / 96
    img_h_in = img.height / 96

    margin = 0.5
    max_w = 13.33 - margin * 2
    max_h = 7.5 - margin * 2
    scale = min(max_w / img_w_in, max_h / img_h_in)
    w, h = img_w_in * scale, img_h_in * scale
    left = (13.33 - w) / 2
    top = (7.5 - h) / 2

    slide.shapes.add_picture(screenshot_path, Inches(left), Inches(top),
                             Inches(w), Inches(h))

    try:
        prs.save(pptx_path)
    except (PermissionError, OSError) as e:
        base, ext = os.path.splitext(pptx_path)
        alt = f"{base}_{int(time.time())}{ext}"
        prs.save(alt)
        print(f"PPTX saved (fallback name): {alt}")
        pptx_path = alt

    try:
        os.remove(screenshot_path)
    except OSError:
        pass

    print(f"PPTX saved: {pptx_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python screenshot.py <html_path> <pptx_path>")
        sys.exit(1)
    main(sys.argv[1], sys.argv[2])
